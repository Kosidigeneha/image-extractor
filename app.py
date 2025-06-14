# app.py
from flask import Flask, request, jsonify, send_file, render_template, send_from_directory
from werkzeug.utils import secure_filename
import os
import fitz  # PyMuPDF for PDF processing
import docx  # python-docx for DOCX processing
import io
import pythoncom 
from PIL import Image
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE
import zipfile
import base64
import tempfile
import shutil
#from docx2pdf import convert
import xml.etree.ElementTree as ET

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
ALLOWED_EXTENSIONS = {'pdf', 'doc', 'docx', 'pptx'}

@app.route('/static/<path:filename>')
def static_files(filename):
    return send_from_directory('static', filename)

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_images_from_pptx(pptx_file):
    """Extract images from PPTX file"""
    images = []
    prs = Presentation(pptx_file)
    
    for slide_number, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                
                try:
                    # Convert image bytes to PIL Image
                    img = Image.open(io.BytesIO(image_bytes))
                    
                    # Skip very small images
                    if img.size[0] < 30 or img.size[1] < 30:
                        continue
                    
                    # Convert to base64
                    img_byte_arr = io.BytesIO()
                    img.save(img_byte_arr, format=img.format if img.format else 'PNG')
                    img_byte_arr = img_byte_arr.getvalue()
                    base64_image = base64.b64encode(img_byte_arr).decode('utf-8')
                    
                    images.append({
                        'data': f"data:image/{img.format.lower() if img.format else 'png'};base64,{base64_image}",
                        'page': slide_number,
                        'original_size': img.size
                    })
                except Exception as e:
                    print(f"Error processing image in slide {slide_number}: {str(e)}")
                    continue
    
    return images

"""def convert_docx_to_pdf(docx_file):

    import win32com.client
    import pythoncom
    
    pythoncom.CoInitialize()
    temp_dir = tempfile.mkdtemp()
    
    try:
        docx_path = os.path.join(temp_dir, 'input.docx')
        pdf_path = os.path.join(temp_dir, 'output.pdf')
        
        with open(docx_path, 'wb') as f:
            f.write(docx_file.read())
        docx_file.seek(0)
        
        try:
            word = win32com.client.DispatchEx('Word.Application')
            doc = word.Documents.Open(docx_path)
            doc.SaveAs(pdf_path, FileFormat=17)
            doc.Close()
            word.Quit()
        except Exception as e:
            print(f"COM automation failed: {str(e)}")
            convert(docx_path, pdf_path)
        
        if os.path.exists(pdf_path):
            pdf_file = open(pdf_path, 'rb')
            return pdf_file, temp_dir
        else:
            raise Exception("PDF conversion failed")
            
    except Exception as e:
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)
        raise Exception(f"Failed to convert DOCX to PDF: {str(e)}")
    finally:
        pythoncom.CoUninitialize()"""

def extract_images_from_pdf(pdf_file, start_page=None, end_page=None):
    """Extract images from PDF file with optional page range"""
    images = []
    pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
    pdf_file.seek(0)
    
    max_pages = len(pdf_document)
    
    if start_page is None or start_page <= 0:
        start_page = 1
    if end_page is None or end_page <= 0:
        end_page = max_pages
    
    start_page = max(1, min(start_page, max_pages))
    end_page = max(start_page, min(end_page, max_pages))
    
    start_idx = start_page - 1
    end_idx = end_page - 1
    
    for page_num in range(start_idx, end_idx + 1):
        page = pdf_document[page_num]
        page_images = []
        
        # Method 1: Extract images using get_images()
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            try:
                base_image = pdf_document.extract_image(xref)
                image_bytes = base_image["image"]
                
                image_format = base_image.get("ext", "").upper()
                if not image_format:
                    image_format = "PNG"
                
                image = Image.open(io.BytesIO(image_bytes))
                
                if image.size[0] < 30 or image.size[1] < 30:
                    continue
                
                img_byte_arr = io.BytesIO()
                image.save(img_byte_arr, format=image_format)
                img_byte_arr = img_byte_arr.getvalue()
                
                base64_image = base64.b64encode(img_byte_arr).decode('utf-8')
                page_images.append({
                    'data': f"data:image/{image_format.lower()};base64,{base64_image}",
                    'page': page_num + 1,
                    'original_size': image.size
                })
            except Exception as e:
                print(f"Error extracting image {img_index} from page {page_num + 1}: {str(e)}")
                continue
        
        # Method 2: Extract embedded images from page pixmap
        try:
            pix = page.get_pixmap()
            if pix.width > 0 and pix.height > 0:
                img_bytes = pix.tobytes()
                image = Image.frombytes("RGB", [pix.width, pix.height], img_bytes)
                
                if image.getbbox():
                    img_byte_arr = io.BytesIO()
                    image.save(img_byte_arr, format='PNG')
                    img_byte_arr = img_byte_arr.getvalue()
                    
                    base64_image = base64.b64encode(img_byte_arr).decode('utf-8')
                    page_images.append({
                        'data': f"data:image/png;base64,{base64_image}",
                        'page': page_num + 1,
                        'original_size': image.size
                    })
        except Exception as e:
            print(f"Error extracting page pixmap from page {page_num + 1}: {str(e)}")
        
        seen_data = set()
        unique_images = []
        for img in page_images:
            if img['data'] not in seen_data:
                seen_data.add(img['data'])
                unique_images.append(img)
        
        images.extend(unique_images)
    
    return images

"""def extract_images_from_doc(doc_file):
    
    pdf_file, temp_dir = convert_docx_to_pdf(doc_file)
    
    try:
        images = extract_images_from_pdf(pdf_file)
        return images
    finally:
        pdf_file.close()
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)"""

def get_document_page_count(file, file_type):
    """Get page count from PDF, DOC, DOCX, or PPTX files"""
    if file_type == 'pdf':
        pdf_document = fitz.open(stream=file.read(), filetype="pdf")
        count = len(pdf_document)
        file.seek(0)
        return count
    elif file_type == 'pptx':
        prs = Presentation(file)
        count = len(prs.slides)
        file.seek(0)
        return count
    elif file_type in ['doc', 'docx']:
        try:
            pdf_file, temp_dir = convert_docx_to_pdf(file)
            pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
            count = len(pdf_document)
            pdf_file.close()
            if temp_dir and os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)
            file.seek(0)
            return count
        except Exception as e:
            raise Exception(f"Error getting page count: {str(e)}")
    return 1

def calculate_grid_positions(images_per_page):
    if images_per_page <= 2:
        return [(0, 0, 1, 1)] if images_per_page == 1 else [(0, 0, 0.5, 1), (0.5, 0, 1, 1)]
    elif images_per_page <= 4:
        return [(0, 0, 0.5, 0.5), (0.5, 0, 1, 0.5),
                (0, 0.5, 0.5, 1), (0.5, 0.5, 1, 1)][:images_per_page]
    else:
        grid = []
        cols = 4 if images_per_page > 6 else 3
        rows = (images_per_page + cols - 1) // cols
        cell_width = 1.0 / cols
        cell_height = 1.0 / rows
        
        for i in range(images_per_page):
            row = i // cols
            col = i % cols
            grid.append((
                col * cell_width,
                row * cell_height,
                (col + 1) * cell_width,
                (row + 1) * cell_height
            ))
        return grid

def create_pdf(images, images_per_page):
    pdf_buffer = io.BytesIO()
    pdf = fitz.open()
    
    images_per_page = min(max(1, int(images_per_page)), 8)
    
    for i in range(0, len(images), images_per_page):
        page = pdf.new_page()
        page_width = page.rect.width
        page_height = page.rect.height
        
        group_images = images[i:i + images_per_page]
        grid = calculate_grid_positions(len(group_images))
        
        margin = 28.35  # 1cm margin
        
        for img_data, position in zip(group_images, grid):
            try:
                left = position[0] * (page_width - 2 * margin) + margin
                top = position[1] * (page_height - 2 * margin) + margin
                width = (position[2] - position[0]) * (page_width - 2 * margin)
                height = (position[3] - position[1]) * (page_height - 2 * margin)
                
                img_data_str = img_data['data'] if isinstance(img_data, dict) else img_data
                img_bytes = base64.b64decode(img_data_str.split(',')[1])
                
                rect = fitz.Rect(left, top, left + width - margin/2, top + height - margin/2)
                page.insert_image(rect, stream=img_bytes)
                
                if isinstance(img_data, dict) and 'page' in img_data:
                    page.insert_text((left, top + height + 5),
                                   f"Page {img_data['page']}",
                                   fontsize=8)
            except Exception as e:
                continue
    
    pdf.save(pdf_buffer)
    pdf_buffer.seek(0)
    return pdf_buffer

def create_pptx(images, images_per_slide):
    prs = Presentation()
    images_per_slide = min(max(1, int(images_per_slide)), 8)
    margin = Cm(1)
    
    for i in range(0, len(images), images_per_slide):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        available_width = prs.slide_width - (2 * margin)
        available_height = prs.slide_height - (2 * margin)
        
        group_images = images[i:i + images_per_slide]
        grid = calculate_grid_positions(len(group_images))
        
        for img_data, position in zip(group_images, grid):
            try:
                left = margin + (position[0] * available_width)
                top = margin + (position[1] * available_height)
                width = (position[2] - position[0]) * available_width - margin/2
                height = (position[3] - position[1]) * available_height - margin/2
                
                img_data_str = img_data['data'] if isinstance(img_data, dict) else img_data
                img_bytes = base64.b64decode(img_data_str.split(',')[1])
                img_io = io.BytesIO(img_bytes)
                
                pic = slide.shapes.add_picture(img_io, left, top, width, height)
                
                if pic.width / pic.height > width / height:
                    new_width = width
                    new_height = width * pic.height / pic.width
                    pic.width = new_width
                    pic.height = new_height
                    pic.top = top + (height - new_height) / 2
                else:
                    new_height = height
                    new_width = height * pic.width / pic.height
                    pic.height = new_height
                    pic.width = new_width
                    pic.left = left + (width - new_width) / 2
                
                if isinstance(img_data, dict) and 'page' in img_data:
                    tx_box = slide.shapes.add_textbox(left, top + height + Cm(0.2), 
                                                    width, Cm(0.6))
                    tx_box.text_frame.text = f"Page {img_data['page']}"
                    tx_box.text_frame.paragraphs[0].alignment = 1
            except Exception as e:
                continue
    
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer

def create_zip(images):
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
        for i, img_data in enumerate(images):
            try:
                img_data_str = img_data['data'] if isinstance(img_data, dict) else img_data
                img_bytes = base64.b64decode(img_data_str.split(',')[1])
                
                page_info = f"_page{img_data['page']}" if isinstance(img_data, dict) and 'page' in img_data else ""
                filename = f"image_{i+1}{page_info}.png"
                
                zip_file.writestr(filename, img_bytes)
            except Exception as e:
                continue
    
    zip_buffer.seek(0)
    return zip_buffer

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/api/get-page-count', methods=['POST'])
def get_page_count():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type'}), 400
    
    try:
        file_type = file.filename.rsplit('.', 1)[1].lower()
        page_count = get_document_page_count(file, file_type)
        return jsonify({'pageCount': page_count})
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/extract-images', methods=['POST'])
def extract_images():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type'}), 400
    
    temp_dir = None
    converted_pdf = None
    
    try:
        # Get page range from form data
        start_page = request.form.get('startPage')
        end_page = request.form.get('endPage')
        
        # Convert to int only if values are provided
        start_page = int(start_page) if start_page and start_page.strip() else None
        end_page = int(end_page) if end_page and end_page.strip() else None
        
        file_type = file.filename.rsplit('.', 1)[1].lower()
        
        if file_type in ['doc', 'docx']:
            # Convert DOCX to PDF first
            converted_pdf, temp_dir = convert_docx_to_pdf(file)
            images = extract_images_from_pdf(converted_pdf, start_page, end_page)
        else:  # pdf
            images = extract_images_from_pdf(file, start_page, end_page)
        
        if not images:
            return jsonify({'images': []})
        
        return jsonify({'images': images})
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    finally:
        # Clean up temporary files
        if converted_pdf:
            converted_pdf.close()
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)

@app.route('/api/convert', methods=['POST'])
def convert_images():
    data = request.get_json()
    if not data or 'images' not in data or 'format' not in data:
        return jsonify({'error': 'Invalid request data'}), 400
    
    images = data['images']
    output_format = data['format']
    images_per_page = min(max(1, int(data.get('imagesPerPage', 1))), 8)
    
    try:
        if output_format == 'pdf':
            pdf_buffer = create_pdf(images, images_per_page)
            return send_file(pdf_buffer, mimetype='application/pdf',
                           as_attachment=True, download_name='extracted_images.pdf')
        
        elif output_format == 'pptx':
            pptx_buffer = create_pptx(images, images_per_page)
            return send_file(pptx_buffer,
                           mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
                           as_attachment=True, download_name='extracted_images.pptx')
        
        elif output_format == 'zip':
            zip_buffer = create_zip(images)
            return send_file(zip_buffer,
                           mimetype='application/zip',
                           as_attachment=True, download_name='extracted_images.zip')
        
        else:
            return jsonify({'error': 'Invalid output format'}), 400
            
    except Exception as e:
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5000)
